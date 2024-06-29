from transformers import AutoTokenizer, AutoModel
import torch
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import nltk
import os
import pickle
from PyQt6.QtCore import QObject, pyqtSignal
from docx import Document
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
import chardet

# Download Punkt sentence tokenizer
nltk.download('punkt', quiet=False)

# ---- Embedding Model ----
def meanpooling(output, mask):
    embeddings = output[0]
    mask = mask.unsqueeze(-1).expand(embeddings.size()).float()
    return torch.sum(embeddings * mask, 1) / torch.clamp(mask.sum(1), min=1e-9)

tokenizer = AutoTokenizer.from_pretrained("neuml/pubmedbert-base-embeddings")
model = AutoModel.from_pretrained("neuml/pubmedbert-base-embeddings")

device = torch.device("cpu") # You can change it to 'cuda' if you have a recent Nvidia GPU with at least 6GB of VRAM.
model.to(device)

def get_embedding(texts, batch_size=32):
    """Calculates embeddings for a list of texts in batches."""
    all_embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i : i + batch_size]
        inputs = tokenizer(batch, padding=True, truncation=False, return_tensors='pt').to(device)
        with torch.no_grad():
            output = model(**inputs)
        batch_embeddings = meanpooling(output, inputs['attention_mask']).to(device).numpy()
        all_embeddings.append(batch_embeddings)
    return np.concatenate(all_embeddings, axis=0)

# ---- Search Engine ----
class LocalSemanticSearchEngine(QObject):
    progress_update = pyqtSignal(int)

    def __init__(self, embeddings_file="embeddings.pkl"):
        super().__init__()
        self.corpus = {}
        self.embeddings = {}
        self.embeddings_file = embeddings_file

    def load_documents(self, path):
        if os.path.isdir(path):
            self._load_embeddings_from_disk(path)
            files = [f for f in os.listdir(path) if f.lower().endswith((".txt", ".pdf", ".docx", ".pptx"))]
            total_files = len(files)
            for i, filename in enumerate(files):
                file_path = os.path.join(path, filename)
                if file_path not in self.embeddings:
                    try:
                        self.load_document(file_path)
                    except Exception as e:
                        print(f"Error loading document {filename}: {str(e)}")
                self.progress_update.emit(int((i + 1) / total_files * 100))
            self._save_embeddings_to_disk(path)
        elif os.path.isfile(path) and path.lower().endswith((".txt", ".pdf", ".docx", ".pptx")):
            self.load_document(path)
            self._save_embeddings_to_disk(os.path.dirname(path))
            self.progress_update.emit(100)
        else:
            raise ValueError("Invalid path provided. Please provide a supported file or a directory.")

    def load_document(self, file_path):
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == ".txt":
            text = self.read_text_file(file_path)
        elif ext == ".pdf":
            text = pdf_extract_text(file_path)
        elif ext == ".docx":
            text = self.extract_text_from_docx(file_path)
        elif ext == ".pptx":
            text = self.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        if not text.strip():
            print(f"Warning: No text content found in {file_path}")
            return

        sentences = nltk.sent_tokenize(text)
        if not sentences:
            print(f"Warning: No sentences extracted from {file_path}")
            return

        embeddings = get_embedding(sentences)
        if embeddings.size == 0:
            print(f"Warning: No embeddings generated for {file_path}")
            return

        self.corpus[file_path] = sentences
        self.embeddings[file_path] = embeddings

    def read_text_file(self, file_path):
        with open(file_path, 'rb') as file:
            raw_data = file.read()

        result = chardet.detect(raw_data)
        encoding = result['encoding']

        try:
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
            return content
        except Exception as e:
            print(f"Error reading file: {e}")
            return ""

    def extract_text_from_docx(self, file_path):
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _load_embeddings_from_disk(self, path):
        embeddings_path = os.path.join(path, self.embeddings_file)
        if os.path.exists(embeddings_path):
            with open(embeddings_path, "rb") as f:
                self.embeddings = pickle.load(f)
                self.corpus = {file_path: nltk.sent_tokenize(open(file_path, 'r', encoding='utf-8').read())
                               for file_path in self.embeddings.keys()}

    def _save_embeddings_to_disk(self, path):
        embeddings_path = os.path.join(path, self.embeddings_file)
        with open(embeddings_path, "wb") as f:
            pickle.dump(self.embeddings, f)

    def search(self, query):
        query_embedding = get_embedding([query])[0]
        max_similarity = -1
        most_similar_file = None

        for file_path, file_embeddings in self.embeddings.items():
            similarities = cosine_similarity(query_embedding.reshape(1, -1), file_embeddings)
            max_sim = np.max(similarities)
            if max_sim > max_similarity:
                max_similarity = max_sim
                most_similar_file = file_path

        return most_similar_file, max_similarity

# ---- RAG Enhanced LLM Inference ----
class RAGEnabledLLM:
    def __init__(self, llm, search_engine, max_token_length):
        self.llm = llm
        self.search_engine = search_engine
        self.max_token_length = int(max_token_length)
        self.response = ""
        self.max_context_length = 1024

    def generate_response(self, query, context_length, similarity_threshold=0.4):
        most_similar_file, similarity_score = self.search_engine.search(query)

        if most_similar_file and similarity_score > similarity_threshold:
            with open(most_similar_file, 'r', encoding='utf-8') as file:
                context = file.read()
            context = context[:min(
                int(context_length), self.max_context_length - len(query) - 100)]
        else:
            context = ""

        prompt = f"""
        You are a helpful and informative AI assistant. 
        Context: '{context}'

        Question: {query}

        Answer the question based on the context and your prior knowledge. 
        If you don't have enough information to answer the question, say that.
        After your answer, provide 3 follow-up questions that are relevant to the context and the user's original question.

        Format your response like this:

        ```json
        {{{{
         "answer": "Your answer here.",
         "follow_up_questions": [
          "Question 1",
          "Question 2",
          "Question 3"
         ]
        }}}}
        ```
        """

        # Ensure context_length is an integer:
        if not isinstance(context_length, int):
            raise TypeError(f"context_length must be an integer, got {type(context_length)}")

        for token in self.llm(prompt, max_tokens=self.max_token_length, stop=["<|end_of_text|>"], echo=False, stream=True):
            yield token['choices'][0]['text']

    def generate_response_without_rag(self, query, context_length):
        prompt = f"""
        You are a helpful and informative AI assistant.

        Question: {query}

        Answer the question based on your own knowledge.
        If you don't have enough information to answer the question, say that.
        After your answer, provide 3 follow-up questions that are relevant to the user's original question.

        Format your response like this:

        ```json
        {{{{
         "answer": "Your answer here.",
         "follow_up_questions": [
          "Question 1",
          "Question 2",
          "Question 3"
         ]
        }}}}
        ```
        """

        # Ensure context_length is an integer:
        if not isinstance(context_length, int):
            raise TypeError(f"context_length must be an integer, got {type(context_length)}")

        for token in self.llm(prompt, max_tokens=self.max_token_length, stop=["<|end_of_text|>"], echo=False, stream=True):
            yield token['choices'][0]['text']